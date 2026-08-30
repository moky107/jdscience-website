"""Original BTEC Unit 1 Chemistry lessons."""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from . import diagrams as dg
from .lessons_extra import (
    add_atomic_depth,
    add_covalent_depth,
    add_electron_depth,
    add_ionic_depth,
    add_metallic_depth,
)
from .theme import (
    CREAM,
    TEAL,
    add_text,
    answer_slide,
    blank,
    content_slide,
    diagram_slide,
    heading,
    misconception_slide,
    new_presentation,
    objectives_slide,
    question_slide,
    save_prs,
    section_slide,
    activity_slide,
    table_slide,
    title_slide,
    two_col,
)


def build_atomic_structure(out):
    d = dg.all_diagrams()
    prs = new_presentation()
    title = "BTEC Unit 1 Chemistry: Atomic Structure"
    title_slide(
        prs,
        title,
        "Atoms, subatomic particles, atomic number, mass number and isotopes — taught as an original Unit 1 Chemistry lesson.",
        "BTEC Level 3  ·  Applied Science  ·  Unit 1  ·  Chemistry",
        d["atom_labelled"],
    )
    objectives_slide(prs, [
        "State the relative charge, relative mass and location of protons, neutrons and electrons.",
        "Use atomic number and mass number to calculate numbers of subatomic particles.",
        "Explain what isotopes are and why they have the same chemical properties.",
        "Calculate a relative atomic mass from isotope abundances.",
        "Describe how ions form when atoms lose or gain electrons.",
    ])
    content_slide(prs, "Where this sits in Unit 1", [
        "Unit 1 assesses Principles and Applications of Science I.",
        "Chemistry paper time is typically about 40 minutes within the 90-minute exam.",
        "Atomic structure underpins bonding, the periodic table and amount-of-substance work.",
        "You will reuse proton number, electron arrangement and ions in later lessons.",
        "Keep units, symbols and definitions precise — examiners reward exact wording.",
    ], ["Have a periodic table available.", "Show working in every calculation."], "Teacher note")
    question_slide(prs, "Prior-knowledge retrieval", [
        "Name the three subatomic particles found in an atom.",
        "Which particle determines which element an atom is?",
        "What is the difference between an atom and a molecule?",
        "Is the nucleus positively charged, negatively charged, or neutral?",
    ], "Do now")
    answer_slide(prs, "Retrieval answers", [
        "Proton, neutron and electron.",
        "The proton number (atomic number).",
        "An atom is a single particle of an element; a molecule contains two or more atoms chemically bonded.",
        "The nucleus is positively charged because it contains protons.",
    ])
    section_slide(prs, "The nuclear model of the atom", "A small, dense nucleus surrounded by electrons in shells")
    diagram_slide(prs, "A labelled atom", d["atom_labelled"], "Original schematic — not to scale.", [
        "Almost all of the mass is in the nucleus.",
        "The nucleus occupies a tiny fraction of the atom's volume.",
        "Electrons occupy shells (energy levels) around the nucleus.",
        "A neutral atom has equal numbers of protons and electrons.",
    ])
    content_slide(prs, "Why the nuclear model matters", [
        "Chemical behaviour depends on the outer electrons, not the nucleus.",
        "Nuclear changes (radioactivity) are a different type of process from chemical reactions.",
        "Mass number and atomic number let us identify an isotope uniquely.",
        "Ions form when the electron count changes; the nucleus stays the same.",
    ])
    table_slide(
        prs,
        "Subatomic particles",
        ["Particle", "Relative charge", "Relative mass", "Where found"],
        [
            ["Proton", "+1", "1", "Nucleus"],
            ["Neutron", "0", "1", "Nucleus"],
            ["Electron", "−1", "Very small (≈ 1/1836)", "Shells"],
        ],
    )
    content_slide(prs, "Protons", [
        "The number of protons is the atomic number, Z.",
        "Every atom of a given element has the same number of protons.",
        "Change the proton number and you change the element.",
        "Protons give the nucleus its positive charge.",
        "In a neutral atom, electron number equals proton number.",
    ], ["Z is always a whole number.", "It is shown at the bottom left of a nuclide symbol."], "Symbol")
    content_slide(prs, "Neutrons", [
        "Neutrons are uncharged particles in the nucleus.",
        "They add mass without changing the element.",
        "Isotopes of the same element differ in neutron number.",
        "Neutron number = mass number − atomic number.",
        "Too few or too many neutrons can make a nucleus unstable (radioisotopes).",
    ])
    content_slide(prs, "Electrons", [
        "Electrons have a relative charge of −1 and a very small mass.",
        "They occupy shells; the first shell holds a maximum of 2 electrons.",
        "The second and third shells are treated as holding up to 8 electrons at this level.",
        "Chemical bonding is about rearranging outer-shell electrons.",
        "Removing an electron from a gaseous atom requires ionisation energy.",
    ])
    diagram_slide(prs, "Reading a nuclide symbol", d["sodium_nuclide"], "Mass number top left, atomic number bottom left.", [
        "General form: mass number A over atomic number Z, then the element symbol.",
        "A = protons + neutrons.",
        "Z = protons.",
        "Electrons in a neutral atom = Z.",
        "For an ion, adjust electrons by the charge.",
    ])
    content_slide(prs, "Worked example — sodium-23", [
        "23Na has A = 23 and Z = 11.",
        "Protons = 11.",
        "Neutrons = 23 − 11 = 12.",
        "Electrons in the neutral atom = 11.",
        "The Na+ ion has 10 electrons because one electron has been removed.",
    ])
    question_slide(prs, "Calculate p, n and e", [
        "12C  (neutral atom)",
        "37Cl (neutral atom)",
        "27Al3+ ",
        "16O2− ",
    ], "Worked practice")
    answer_slide(prs, "Particle-count answers", [
        "12C: 6 p, 6 n, 6 e.",
        "37Cl: 17 p, 20 n, 17 e.",
        "27Al3+: 13 p, 14 n, 10 e (13 − 3).",
        "16O2−: 8 p, 8 n, 10 e (8 + 2).",
    ])
    section_slide(prs, "Isotopes", "Same element, different mass number")
    content_slide(prs, "Defining isotopes", [
        "Isotopes are atoms of the same element with the same proton number but different neutron numbers.",
        "They therefore have different mass numbers.",
        "Chemical properties are almost identical because electron arrangement is the same.",
        "Physical properties that depend on mass (density, rate of diffusion) can differ.",
        "Some isotopes are stable; others are radioactive.",
    ], ["Do not say isotopes have different proton numbers.", "That would make them different elements."], "Watch this")
    two_col(
        prs,
        "Two chlorine isotopes",
        "Chlorine-35",
        ["17 protons", "18 neutrons", "17 electrons (atom)", "More abundant in nature"],
        "Chlorine-37",
        ["17 protons", "20 neutrons", "17 electrons (atom)", "Less abundant in nature"],
    )
    content_slide(prs, "Relative atomic mass, Ar", [
        "Most elements exist as a mixture of isotopes.",
        "Ar is the weighted mean mass of an atom of an element compared with 1/12 of the mass of a 12C atom.",
        "Use percentage abundances from data, not from memory, in an exam.",
        "Ar has no units.",
        "It is the value shown on most periodic tables.",
    ])
    content_slide(prs, "Worked example — chlorine Ar", [
        "A sample is 75.0% 35Cl and 25.0% 37Cl.",
        "Ar = (75.0 × 35 + 25.0 × 37) / 100",
        "Ar = (2625 + 925) / 100",
        "Ar = 3550 / 100 = 35.5",
        "This matches the familiar periodic-table value.",
    ])
    question_slide(prs, "Relative atomic mass practice", [
        "Boron is 20.0% 10B and 80.0% 11B. Calculate Ar.",
        "A copper sample is 69.0% 63Cu and 31.0% 65Cu. Calculate Ar to 1 d.p.",
        "Why is Ar for chlorine not a whole number?",
    ])
    answer_slide(prs, "Ar answers", [
        "Ar(B) = (20×10 + 80×11)/100 = 10.8",
        "Ar(Cu) = (69×63 + 31×65)/100 = 63.6",
        "Chlorine is a mixture of 35Cl and 37Cl, so the weighted mean is not an integer.",
    ])
    section_slide(prs, "Ions from atoms", "Changing electrons, not protons")
    content_slide(prs, "How ions form", [
        "A positive ion (cation) forms when an atom loses one or more electrons.",
        "A negative ion (anion) forms when an atom gains one or more electrons.",
        "The proton number is unchanged, so it is still the same element.",
        "Metals typically form cations; non-metals typically form anions.",
        "The charge equals protons minus electrons.",
    ])
    activity_slide(prs, "Build four nuclide cards", [
        "On mini-whiteboards, draw cards for 24Mg, 24Mg2+, 19F and 19F−.",
        "On each card show A, Z, protons, neutrons and electrons.",
        "Swap with a partner and mark one correction if needed.",
        "Be ready to explain how the 2+ and − ions were obtained.",
    ], "10 minutes")
    misconception_slide(prs, [
        ("The mass number is the number of neutrons.", "Mass number is protons plus neutrons."),
        ("Electrons sit inside the nucleus with the protons.", "Electrons occupy shells outside the nucleus."),
        ("Isotopes have different numbers of protons.", "Isotopes have the same protons and different neutrons."),
        ("Ions are formed by changing the number of protons.", "Ions form when electrons are lost or gained."),
    ])
    add_atomic_depth(prs)
    question_slide(prs, "Exam-style practice", [
        "A student writes that 14C and 14N are isotopes. Explain why this is incorrect. (2)",
        "An element has two isotopes, 69X (60.0%) and 71X (40.0%). Calculate Ar. (2)",
        "Explain why atoms are electrically neutral even though they contain charged particles. (2)",
    ], "Exam-style practice")
    answer_slide(prs, "Exam-style answers", [
        "Isotopes must have the same proton number / be the same element. 14C has 6 protons; 14N has 7 protons, so they are different elements.",
        "Ar = (60×69 + 40×71)/100 = 69.8",
        "The number of protons equals the number of electrons, so the + and − charges cancel.",
    ])
    content_slide(prs, "Challenge", [
        "A mass spectrum of neon shows peaks at m/z 20, 21 and 22 with relative heights 90.5, 0.3 and 9.2.",
        "Identify which peak is 20Ne.",
        "Estimate Ar and comment on whether 21Ne can be ignored in a 3 s.f. answer.",
        "Explain one reason a fourth peak might appear if the sample were ionised as Ne2+.",
    ])
    content_slide(prs, "Plenary — say it in 30 seconds", [
        "Define atomic number and mass number.",
        "Give one pair of isotopes and state what is the same and what is different.",
        "State the formula for Ar from two isotopes.",
        "Explain in one sentence how a 2− ion forms.",
    ])
    content_slide(prs, "Independent practice", [
        "Complete the matching worksheet for this lesson.",
        "Learn the relative charges and masses of the three particles.",
        "Calculate p, n and e for five nuclides of your choice from a periodic table.",
        "Write a six-mark plan: 'Explain the existence of isotopes and how Ar is calculated.'",
    ], ["Worksheet sold separately as a companion resource.", "Answers are on the teacher sheet."], "Homework")
    save_prs(prs, out, title, "BTEC Unit 1 Chemistry", "Atomic structure")
    return title, len(prs.slides)


def build_electron_configuration(out):
    d = dg.all_diagrams()
    prs = new_presentation()
    title = "BTEC Unit 1 Chemistry: Electron Configuration"
    title_slide(prs, title, "Shells, subshells and how electron arrangement links to the periodic table.", "BTEC Level 3  ·  Applied Science  ·  Unit 1  ·  Chemistry", d["shells_diagram"])
    objectives_slide(prs, [
        "Write electron configurations using shells (2,8,8) for the first 20 elements.",
        "Describe s, p and d subshells and state their capacities.",
        "Construct configurations using the Aufbau order for s- and p-block elements.",
        "Relate outer-shell electrons to group number for main-group elements.",
        "Use electron arrangement to predict simple ion charges.",
    ])
    question_slide(prs, "Prior knowledge", [
        "How many electrons can the first shell hold?",
        "Why do atoms form ions?",
        "What does the group number of sodium tell you?",
        "Which particle is rearranged during bonding?",
    ], "Do now")
    answer_slide(prs, "Prior-knowledge answers", [
        "A maximum of two electrons.",
        "To obtain a more stable (often full) outer shell.",
        "Sodium is in Group 1, so it has one outer-shell electron.",
        "Electrons.",
    ])
    section_slide(prs, "Shells first", "Energy levels around the nucleus")
    diagram_slide(prs, "Electron shells", d["shells_diagram"], "Shell capacity used at BTEC Unit 1: 2, 8, 8 for the first 20 elements.", [
        "Electrons occupy the lowest available shell first.",
        "A full outer shell is associated with chemical stability.",
        "Noble gases already have a full outer shell.",
        "Period number matches the number of occupied shells for main-group atoms.",
    ])
    table_slide(prs, "Shell notation for the first 20 elements (selected)", ["Element", "Z", "Shells", "Ions commonly formed"], [
        ["He", "2", "2", "None (noble gas)"],
        ["O", "8", "2,6", "O2−"],
        ["Na", "11", "2,8,1", "Na+"],
        ["Cl", "17", "2,8,7", "Cl−"],
        ["Ca", "20", "2,8,8,2", "Ca2+"],
    ])
    content_slide(prs, "Worked example — chlorine", [
        "Chlorine has 17 electrons.",
        "Fill the first shell: 2 electrons remain 15.",
        "Fill the second shell: 8 electrons remain 7.",
        "Configuration: 2,8,7.",
        "Chlorine commonly gains one electron to become 2,8,8 (Cl−).",
    ])
    question_slide(prs, "Write the shell configurations", [
        "Nitrogen (Z = 7)",
        "Magnesium (Z = 12)",
        "Phosphorus (Z = 15)",
        "Potassium (Z = 19)",
        "The ion S2− (Z = 16)",
    ])
    answer_slide(prs, "Shell-configuration answers", [
        "N: 2,5",
        "Mg: 2,8,2",
        "P: 2,8,5",
        "K: 2,8,8,1",
        "S2−: 2,8,8  (atom is 2,8,6; two electrons added)",
    ])
    section_slide(prs, "Subshells", "s, p and d")
    content_slide(prs, "Why subshells are needed", [
        "Shells split into subshells of slightly different energy.",
        "This explains the structure of the periodic table more precisely.",
        "Unit 1 expects you to recognise s, p and d blocks.",
        "You should be able to write configurations such as 1s2 2s2 2p6.",
        "Hund's rule and Pauli exclusion are used when drawing orbital boxes.",
    ])
    two_col(prs, "Subshell capacities", "Type and orbitals", [
        "s subshell: 1 orbital",
        "p subshell: 3 orbitals",
        "d subshell: 5 orbitals",
        "Each orbital holds a maximum of 2 electrons with opposite spin.",
    ], "Maximum electrons", [
        "s: 2 electrons",
        "p: 6 electrons",
        "d: 10 electrons",
        "The n = 2 shell is 2s + 2p (8 electrons in total).",
    ])
    content_slide(prs, "Aufbau order used at this level", [
        "Fill 1s, then 2s, then 2p, then 3s, then 3p, then 4s, then 3d.",
        "4s is filled before 3d for isolated atoms of the first row of transition metals.",
        "Write superscripts for electron counts.",
        "Example: oxygen is 1s2 2s2 2p4.",
        "Example: calcium is 1s2 2s2 2p6 3s2 3p6 4s2.",
    ])
    content_slide(prs, "Worked example — phosphorus", [
        "Phosphorus has 15 electrons.",
        "1s2 2s2 2p6 uses 10 electrons; 5 remain.",
        "3s2 3p3 uses the remaining 5.",
        "Full configuration: 1s2 2s2 2p6 3s2 3p3.",
        "Noble-gas shorthand: [Ne] 3s2 3p3.",
    ])
    activity_slide(prs, "Build configurations", [
        "Write full subshell configurations for F, Al, Ar and Ca.",
        "Circle the outer-shell electrons.",
        "State the most likely ion for F and Al and rewrite the ion configuration.",
        "Check a partner's Ca configuration carefully — 4s comes after 3p.",
    ], "12 minutes")
    answer_slide(prs, "Activity answers", [
        "F: 1s2 2s2 2p5 → F− is 1s2 2s2 2p6",
        "Al: 1s2 2s2 2p6 3s2 3p1 → Al3+ is 1s2 2s2 2p6",
        "Ar: 1s2 2s2 2p6 3s2 3p6",
        "Ca: 1s2 2s2 2p6 3s2 3p6 4s2",
    ])
    content_slide(prs, "Link to the periodic table", [
        "Group 1 atoms have one outer s electron (ns1).",
        "Group 2 atoms have ns2.",
        "Group 7 (17) atoms have ns2 np5.",
        "Noble gases have ns2 np6 (except helium, 1s2).",
        "The d-block corresponds to filling 3d (and later d subshells).",
    ], ["Outer electrons control bonding and ion charge.", "Inner electrons are not usually shown in simple bonding diagrams."], "Teaching point")
    misconception_slide(prs, [
        ("The third shell always holds 18 electrons in Unit 1 examples of the first 20 elements.", "For K and Ca we still use 2,8,8,1 and 2,8,8,2 at this level."),
        ("2p4 means four p subshells.", "It means four electrons in the 2p subshell."),
        ("Ions keep the same electron configuration as the atom.", "Ions have lost or gained electrons, so the configuration changes."),
    ])
    add_electron_depth(prs)
    question_slide(prs, "Exam-style practice", [
        "Write the electron configuration of Mg2+ and explain why this ion is stable. (3)",
        "A student writes chlorine as 2,8,8. Identify the error and give the correct configuration of the atom. (2)",
        "Explain, using electron configuration, why oxygen typically forms a 2− ion. (3)",
    ], "Exam-style practice")
    answer_slide(prs, "Exam-style answers", [
        "Mg atom is 2,8,2 / 1s2 2s2 2p6 3s2. Mg2+ is 2,8 / 1s2 2s2 2p6. This is a full outer shell / noble-gas arrangement, so it is stable.",
        "The atom has 17 electrons, not 18. Correct atom configuration is 2,8,7.",
        "Oxygen is 2,6. Gaining two electrons gives 2,8. A full outer shell is more stable, so O2− forms.",
    ])
    content_slide(prs, "Plenary", [
        "On a whiteboard: element, shell configuration, subshell configuration, likely ion.",
        "Use sodium, oxygen and calcium.",
        "Explain one link between group number and outer electrons.",
    ])
    content_slide(prs, "Independent practice", [
        "Complete the electron-configuration worksheet.",
        "Learn configurations for the first 20 elements.",
        "Preview ionic bonding: why Na and Cl form NaCl.",
    ])
    save_prs(prs, out, title, "BTEC Unit 1 Chemistry", "Electron configuration")
    return title, len(prs.slides)


def build_ionic_bonding(out):
    d = dg.all_diagrams()
    prs = new_presentation()
    title = "BTEC Unit 1 Chemistry: Ionic Bonding"
    title_slide(prs, title, "Electron transfer, ions, giant lattices and the properties they explain.", "BTEC Level 3  ·  Applied Science  ·  Unit 1  ·  Chemistry", d["ionic_lattice"])
    objectives_slide(prs, [
        "Describe ionic bonding as electrostatic attraction between oppositely charged ions.",
        "Explain how ions form by electron transfer between metals and non-metals.",
        "Draw simple representations of ion formation for binary ionic compounds.",
        "Describe the giant ionic lattice and use it to explain melting point, electrical conductivity and solubility.",
        "Deduce formulae of simple ionic compounds from ion charges.",
    ])
    question_slide(prs, "Prior knowledge", [
        "What is the charge on a proton and on an electron?",
        "Write the shell configuration of sodium and of chlorine.",
        "Why do Group 1 metals form 1+ ions?",
        "Name one compound you already know that contains ions.",
    ], "Do now")
    answer_slide(prs, "Retrieval answers", [
        "Proton +1; electron −1.",
        "Na 2,8,1 and Cl 2,8,7.",
        "They have one outer electron which is lost to leave a full outer shell.",
        "Examples include sodium chloride, copper sulfate or calcium carbonate.",
    ])
    section_slide(prs, "What ionic bonding is", "Transfer, not sharing")
    content_slide(prs, "A precise definition", [
        "Ionic bonding is the strong electrostatic attraction between oppositely charged ions.",
        "It typically forms between a metal and a non-metal.",
        "The metal atom loses electrons; the non-metal atom gains electrons.",
        "The attraction acts in all directions in a giant lattice.",
        "Do not describe ionic bonding as 'sharing electrons'.",
    ])
    diagram_slide(prs, "Electron transfer: sodium and chlorine", d["ion_transfer"], "Original particle picture of charge imbalance after transfer.", [
        "Na starts 2,8,1 and becomes Na+ 2,8.",
        "Cl starts 2,8,7 and becomes Cl− 2,8,8.",
        "One electron is transferred.",
        "The ions then attract strongly.",
        "The compound formed is NaCl.",
    ])
    content_slide(prs, "Worked example — magnesium oxide", [
        "Mg is 2,8,2 and needs to lose two electrons to reach 2,8.",
        "O is 2,6 and needs to gain two electrons to reach 2,8.",
        "One Mg atom supplies both electrons to one O atom.",
        "Ions formed: Mg2+ and O2−.",
        "Formula: MgO. The lattice contains these ions in a 1:1 ratio.",
    ])
    content_slide(prs, "Worked example — calcium chloride", [
        "Ca is 2,8,8,2 → loses 2 electrons → Ca2+.",
        "Each Cl is 2,8,7 → gains 1 electron → Cl−.",
        "Two Cl atoms are needed for each Ca atom.",
        "Formula: CaCl2.",
        "Check: total positive charge (+2) equals total negative charge (2 × −1).",
    ])
    activity_slide(prs, "Formulae from charges", [
        "Deduce formulae for compounds of: Na and O; Al and Cl; Mg and N; K and SO4 2−.",
        "Show the charge-balance working.",
        "Name each compound.",
    ], "8 minutes")
    answer_slide(prs, "Formula answers", [
        "Na2O — sodium oxide",
        "AlCl3 — aluminium chloride",
        "Mg3N2 — magnesium nitride",
        "K2SO4 — potassium sulfate",
    ])
    section_slide(prs, "The giant ionic lattice", "Structure explains properties")
    diagram_slide(prs, "Giant ionic lattice", d["ionic_lattice"], "Two-dimensional sketch of a three-dimensional lattice.", [
        "Every cation is surrounded by anions, and vice versa.",
        "The attraction is strong and acts throughout the crystal.",
        "A large amount of energy is needed to separate the ions.",
        "This produces high melting and boiling points.",
    ])
    two_col(prs, "Properties of ionic compounds", "Typical observations", [
        "High melting and boiling points",
        "Usually soluble in water",
        "Do not conduct when solid",
        "Do conduct when molten or aqueous",
        "Often brittle crystals",
    ], "Explanation from structure", [
        "Strong attractions throughout the lattice",
        "Water molecules can surround and separate ions",
        "Ions are fixed in place, so charge cannot move",
        "Ions become free to move and carry charge",
        "A shift in the lattice lines up like charges, so the crystal splits",
    ])
    content_slide(prs, "Electrical conductivity in more detail", [
        "Conduction requires mobile charged particles.",
        "In a solid ionic lattice the ions vibrate but cannot move past one another.",
        "When the solid melts, ions can migrate towards electrodes.",
        "When the solid dissolves, hydrated ions move independently.",
        "This is different from metallic conduction, which uses delocalised electrons.",
    ])
    misconception_slide(prs, [
        ("NaCl molecules travel around in salt solution.", "There are no NaCl molecules; there are separate hydrated Na+ and Cl− ions."),
        ("Ionic compounds conduct because electrons jump from ion to ion in the solid.", "The solid does not conduct; mobile ions are required."),
        ("A 2+ ion is formed by gaining two protons.", "Charge changes by losing or gaining electrons."),
    ])
    add_ionic_depth(prs)
    question_slide(prs, "Exam-style practice", [
        "Describe, in terms of electrons, how ionic bonding arises in magnesium oxide. (4)",
        "Explain why sodium chloride has a high melting point but does not conduct electricity as a solid. (4)",
        "Deduce the formula of the compound formed between Al3+ and O2−. Show the charge balance. (2)",
    ], "Exam-style practice")
    answer_slide(prs, "Exam-style answers", [
        "Mg atom loses 2 electrons to form Mg2+. O atom gains 2 electrons to form O2−. Oppositely charged ions attract electrostatically. A giant lattice of Mg2+ and O2− is formed.",
        "High melting point: strong electrostatic attraction between oppositely charged ions throughout the lattice needs a lot of energy to overcome. Solid does not conduct: ions are not free to move.",
        "Need two Al3+ and three O2− so +6 balances −6. Formula Al2O3.",
    ])
    content_slide(prs, "Plenary", [
        "Write a 20-word definition of ionic bonding.",
        "Sketch two ions and mark the attraction.",
        "Give one property and its structural explanation.",
    ])
    content_slide(prs, "Independent practice", [
        "Complete the ionic bonding worksheet.",
        "Practise formulae until charge balance is automatic.",
        "Next lesson: covalent bonding — sharing rather than transfer.",
    ])
    save_prs(prs, out, title, "BTEC Unit 1 Chemistry", "Ionic bonding")
    return title, len(prs.slides)


def build_covalent_bonding(out):
    d = dg.all_diagrams()
    prs = new_presentation()
    title = "BTEC Unit 1 Chemistry: Covalent Bonding"
    title_slide(prs, title, "Shared electron pairs, simple molecules and giant covalent structures.", "BTEC Level 3  ·  Applied Science  ·  Unit 1  ·  Chemistry", d["covalent_pair"])
    objectives_slide(prs, [
        "Describe a covalent bond as a shared pair of electrons.",
        "Explain why non-metal atoms share electrons.",
        "Represent simple molecules using straight-line bonds and state of outer-shell counts.",
        "Distinguish simple molecular structures from giant covalent structures.",
        "Use structure to explain melting point, solubility and electrical conductivity.",
    ])
    question_slide(prs, "Prior knowledge", [
        "How does ionic bonding differ from sharing electrons?",
        "How many electrons does carbon need to obtain a full outer shell?",
        "Name two elements that exist as diatomic molecules.",
        "What is meant by a full outer shell for elements in Period 2?",
    ], "Do now")
    answer_slide(prs, "Retrieval answers", [
        "Ionic bonding is electron transfer and attraction between ions; covalent bonding is sharing.",
        "Four (carbon is 2,4).",
        "Hydrogen, nitrogen, oxygen, fluorine, chlorine, bromine or iodine.",
        "Eight electrons in the outer shell (helium is the exception with two).",
    ])
    section_slide(prs, "The covalent bond", "A shared pair")
    diagram_slide(prs, "A shared pair of electrons", d["covalent_pair"], "Hydrogen chloride as a simple original schematic.", [
        "Each atom contributes one electron to the shared pair.",
        "The pair is attracted to both nuclei.",
        "That attraction is the covalent bond.",
        "Both atoms obtain a more stable outer shell.",
    ])
    content_slide(prs, "Single, double and triple bonds", [
        "A single bond is one shared pair (H–H, Cl–Cl, C–C).",
        "A double bond is two shared pairs (O=O, C=O).",
        "A triple bond is three shared pairs (N≡N).",
        "Multiple bonds are shorter and stronger than the corresponding single bond between the same atoms.",
        "Carbon regularly forms four covalent bonds in stable molecules.",
    ])
    two_col(prs, "Worked molecules", "Water, H2O", [
        "Oxygen is 2,6 and needs two more electrons.",
        "Each hydrogen shares its one electron.",
        "Oxygen forms two single bonds.",
        "The molecule is bent; this matters later for polarity.",
    ], "Carbon dioxide, CO2", [
        "Carbon needs four electrons; each oxygen needs two.",
        "Carbon forms two double bonds, O=C=O.",
        "The molecule is linear.",
        "It is a simple molecule with weak forces between molecules.",
    ])
    content_slide(prs, "Methane and ammonia", [
        "Methane, CH4: carbon forms four C–H single bonds. Tetrahedral shape.",
        "Ammonia, NH3: nitrogen forms three N–H bonds and has one lone pair.",
        "A lone pair is a pair of outer electrons not used in bonding.",
        "Lone pairs affect shape and can accept a proton (to form NH4+).",
        "These are still simple molecules, not giant lattices.",
    ])
    activity_slide(prs, "Bond counts", [
        "For F2, O2, N2, CH4 and CO2, state the number of shared pairs in the molecule.",
        "Identify any lone pairs on the central atom where relevant.",
        "Predict whether each substance is a gas at room temperature and justify using structure, not memory alone.",
    ], "10 minutes")
    answer_slide(prs, "Bond-count answers", [
        "F2: 1 pair. O2: 2 pairs. N2: 3 pairs.",
        "CH4: 4 pairs; no lone pair on C. CO2: 4 pairs in total as two double bonds; no lone pair on C.",
        "All are simple molecules, so intermolecular forces are weak and they are gases at room temperature.",
    ])
    section_slide(prs, "Two covalent structure types", "Simple molecular versus giant covalent")
    two_col(prs, "Structure types", "Simple molecular", [
        "Strong covalent bonds within molecules",
        "Weak forces between molecules",
        "Low melting and boiling points",
        "Usually do not conduct electricity",
        "Examples: H2O, CO2, I2, CH4",
    ], "Giant covalent", [
        "Vast network of covalent bonds",
        "No separate small molecules",
        "Very high melting points",
        "Usually do not conduct (graphite is the important exception)",
        "Examples: diamond, graphite, silicon dioxide",
    ])
    content_slide(prs, "Diamond and graphite — contrast", [
        "In diamond, each carbon is bonded to four others in a tetrahedral network.",
        "Diamond is extremely hard and does not conduct; all outer electrons are in bonds.",
        "In graphite, each carbon is bonded to three others in layers.",
        "Spare electrons are delocalised between layers, so graphite conducts.",
        "Layers can slide, so graphite is soft and is used as a lubricant and in pencils.",
    ])
    content_slide(prs, "Electrical conductivity of covalent substances", [
        "Simple molecules have no mobile ions or delocalised electrons, so they do not conduct.",
        "Diamond has no delocalised electrons.",
        "Graphite and graphene have delocalised electrons and do conduct.",
        "Molten ionic compounds conduct for a different reason — mobile ions.",
        "Always name the mobile charged particle in an explanation.",
    ])
    misconception_slide(prs, [
        ("Breaking covalent bonds is what happens when water boils.", "Boiling overcomes intermolecular forces; the H–O bonds remain."),
        ("All covalent substances have low melting points.", "Giant covalent structures have very high melting points."),
        ("Covalent bonds form between a metal and a non-metal.", "That pairing usually gives ionic bonding."),
    ])
    add_covalent_depth(prs)
    question_slide(prs, "Exam-style practice", [
        "Describe how a covalent bond forms in a chlorine molecule. (3)",
        "Explain why iodine melts at a much lower temperature than diamond, even though both contain covalent bonds. (4)",
        "Explain why graphite conducts electricity but diamond does not. (3)",
    ], "Exam-style practice")
    answer_slide(prs, "Exam-style answers", [
        "Each chlorine atom has 7 outer electrons. They share one pair. Both atoms then have 8 outer electrons. The shared pair is attracted to both nuclei.",
        "Iodine is simple molecular: melting overcomes weak forces between I2 molecules. Diamond is giant covalent: melting requires breaking many strong C–C bonds, which needs far more energy.",
        "Graphite has delocalised electrons that can move between layers. In diamond every outer electron is held in a C–C bond, so there are no mobile charged particles.",
    ])
    content_slide(prs, "Plenary", [
        "Define a covalent bond in one sentence.",
        "Give one simple-molecular and one giant-covalent example.",
        "State the mobile particle (if any) that allows graphite to conduct.",
    ])
    content_slide(prs, "Independent practice", [
        "Complete the covalent bonding worksheet.",
        "Draw displayed formulae for H2O, NH3, CH4 and CO2 from memory.",
        "Next lesson: metallic bonding and why metals conduct as solids.",
    ])
    save_prs(prs, out, title, "BTEC Unit 1 Chemistry", "Covalent bonding")
    return title, len(prs.slides)


def build_metallic_bonding(out):
    d = dg.all_diagrams()
    prs = new_presentation()
    title = "BTEC Unit 1 Chemistry: Metallic Bonding"
    title_slide(prs, title, "Positive ions, delocalised electrons and the properties of metals.", "BTEC Level 3  ·  Applied Science  ·  Unit 1  ·  Chemistry", d["metallic_lattice"])
    objectives_slide(prs, [
        "Describe metallic bonding as the attraction between positive metal ions and a sea of delocalised electrons.",
        "Explain why metals conduct electricity and heat.",
        "Explain malleability and ductility using the metallic lattice.",
        "Suggest why different metals have different melting points and strengths.",
        "Compare metallic, ionic and covalent structures.",
    ])
    question_slide(prs, "Prior knowledge", [
        "Are metal atoms more likely to lose or gain electrons?",
        "Why does solid sodium chloride not conduct electricity?",
        "Name two physical properties common to most metals.",
        "What does delocalised mean in everyday language?",
    ], "Do now")
    answer_slide(prs, "Retrieval answers", [
        "Lose electrons, forming positive ions.",
        "Its ions are not free to move.",
        "Shiny, conduct heat and electricity, malleable, ductile, high melting points (most).",
        "Not fixed in one place; free to move through the structure.",
    ])
    section_slide(prs, "The metallic lattice", "Ions plus a sea of electrons")
    diagram_slide(prs, "Metallic structure", d["metallic_lattice"], "Original schematic of cations in a delocalised electron cloud.", [
        "Metal atoms lose their outer electrons into a shared cloud.",
        "The remaining particles are positive ions packed closely.",
        "Delocalised electrons attract every ion and bind the lattice.",
        "The attraction is strong and non-directional.",
    ])
    content_slide(prs, "A careful definition", [
        "Metallic bonding is the electrostatic attraction between positive metal ions and delocalised electrons.",
        "It is not the same as ionic bonding: there are no anions.",
        "It is not covalent bonding: electrons are not shared in localised pairs between two atoms.",
        "The number of delocalised electrons per atom is often equal to the number of outer-shell electrons.",
        "More delocalised electrons and higher ionic charge generally strengthen the bond.",
    ])
    two_col(prs, "Explaining metallic properties", "Electrical conductivity", [
        "Delocalised electrons are mobile charged particles.",
        "An applied potential difference makes them drift.",
        "The solid therefore conducts.",
        "Molten metals also conduct.",
    ], "Thermal conductivity", [
        "Electrons transfer kinetic energy rapidly through the lattice.",
        "Ions also vibrate and pass energy to neighbours.",
        "This is why a metal spoon heats quickly in a hot drink.",
        "Name electrons in a full explanation.",
    ])
    content_slide(prs, "Malleability and ductility", [
        "Layers of positive ions can slide over one another when a force is applied.",
        "The delocalised electrons continue to hold the ions together after the slide.",
        "The metal therefore changes shape instead of shattering.",
        "Malleable: can be hammered into sheets. Ductile: can be drawn into wires.",
        "Ionic crystals shatter because sliding lines up like charges.",
    ])
    content_slide(prs, "Strength and melting point", [
        "Group 1 metals are relatively soft and have lower melting points.",
        "They contribute only one delocalised electron per atom and form 1+ ions.",
        "Transition metals often have higher melting points and are harder.",
        "They can contribute more electrons and have stronger metallic bonding.",
        "Alloys are harder than pure metals because different-sized atoms prevent layers sliding easily.",
    ])
    activity_slide(prs, "Three-structure comparison", [
        "Draw a three-column table: ionic, simple covalent, metallic.",
        "For each, state the particles present, the bonding, melting point (high/low), and whether the solid conducts.",
        "Add one example to each column.",
        "Write one sentence explaining the conductivity difference between copper and solid sodium chloride.",
    ], "12 minutes")
    answer_slide(prs, "Comparison checkpoints", [
        "Ionic: ions; electrostatic attraction; high m.p.; solid does not conduct.",
        "Simple covalent: molecules; covalent bonds inside, weak forces between; low m.p.; does not conduct.",
        "Metallic: cations + delocalised electrons; metallic bonding; usually high m.p.; solid does conduct.",
        "Copper has mobile electrons; solid NaCl does not have mobile ions.",
    ])
    misconception_slide(prs, [
        ("Metals conduct because positive ions flow through the wire.", "The ions stay in the lattice; electrons move."),
        ("Metallic bonding is a type of covalent bonding.", "Electrons are delocalised, not shared as localised pairs."),
        ("All metals have very high melting points.", "Mercury is liquid at room temperature; Group 1 metals melt relatively easily."),
    ])
    add_metallic_depth(prs)
    question_slide(prs, "Exam-style practice", [
        "Describe the bonding in solid magnesium. (3)",
        "Explain why magnesium conducts electricity as a solid but magnesium oxide does not. (4)",
        "Explain why metals can be drawn into wires. (3)",
    ], "Exam-style practice")
    answer_slide(prs, "Exam-style answers", [
        "Magnesium consists of Mg2+ ions packed in a lattice. Outer electrons are delocalised. Electrostatic attraction between ions and electrons is metallic bonding.",
        "Mg has delocalised electrons that can move and carry charge. MgO is ionic; in the solid the ions are not free to move, so it does not conduct.",
        "Layers of ions can slide. Delocalised electrons continue to attract the ions, so the metal does not break.",
    ])
    content_slide(prs, "Plenary", [
        "Complete: 'Metallic bonding is the attraction between … and …'",
        "Give one property explained by sliding layers and one explained by mobile electrons.",
        "State one difference between an alloy and a pure metal in terms of structure.",
    ])
    content_slide(prs, "Independent practice", [
        "Complete the metallic bonding worksheet.",
        "Revise all three bonding types together.",
        "Preview intermolecular forces: why CO2 is a gas even though it has strong C=O bonds.",
    ])
    save_prs(prs, out, title, "BTEC Unit 1 Chemistry", "Metallic bonding")
    return title, len(prs.slides)
