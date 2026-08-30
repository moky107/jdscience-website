"""JDScience lesson theme: teal identity, topic titles, discreet footer, varied layouts."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .notation import add_markup_runs, apply_run

ROOT = Path(__file__).resolve().parents[2]
MEDIA = ROOT / "content/lessons/_media"
MEDIA.mkdir(parents=True, exist_ok=True)
LOGO_PATH = MEDIA / "jdscience-footer-logo.png"

W, H = Inches(13.333), Inches(7.5)
FOOTER_H = Inches(0.46)

TEAL = RGBColor(0x00, 0x96, 0x88)
TEAL_DARK = RGBColor(0x00, 0x4D, 0x40)
TEAL_MID = RGBColor(0x00, 0x79, 0x6B)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF0, 0xFD, 0xFA)
CARD = RGBColor(0xF8, 0xFA, 0xFC)
LINE = RGBColor(0xCC, 0xFB, 0xF1)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG = RGBColor(0xFF, 0xF7, 0xED)
ROSE = RGBColor(0x9F, 0x12, 0x39)
ROSE_BG = RGBColor(0xFF, 0xF1, 0xF2)
GREEN = RGBColor(0x04, 0x78, 0x57)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_BG = RGBColor(0xE8, 0xEE, 0xF5)
PURPLE = RGBColor(0x5B, 0x21, 0xB6)
PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)
SKY = RGBColor(0x0E, 0x74, 0x90)
SKY_BG = RGBColor(0xEC, 0xFE, 0xFF)


def ensure_logo() -> Path:
    if LOGO_PATH.exists():
        return LOGO_PATH
    img = Image.new("RGBA", (160, 160), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((0, 0, 159, 159), radius=28, fill=(0, 150, 136, 255))
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 58)
    except OSError:
        font = ImageFont.load_default()
    d.text((80, 78), "JD", fill="white", font=font, anchor="mm")
    img.save(LOGO_PATH)
    return LOGO_PATH


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def add_rect(slide, l, t, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_round(slide, l, t, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.25)
    return shape


def set_run(run, size=18, bold=False, color=INK, name="Calibri", italic=False):
    apply_run(run, size=size, bold=bold, color=color, italic=italic, name=name)


def add_text(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    add_markup_runs(p, text, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, l, t, w, h, items, size=18, color=INK, spacing=10):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = "•  "
        set_run(run, size=size, color=color)
        add_markup_runs(p, item, size=size, color=color)
    return box


def footer(slide, page: int, total: int, kicker: str = ""):
    add_rect(slide, 0, H - FOOTER_H, W, FOOTER_H, TEAL_DARK)
    logo = str(ensure_logo())
    slide.shapes.add_picture(logo, Inches(0.22), H - FOOTER_H + Inches(0.07), Inches(0.32), Inches(0.32))
    add_text(slide, Inches(0.62), H - FOOTER_H + Inches(0.08), Inches(2.4), Inches(0.30), "jdscience.co.uk", 11, True, WHITE)
    if kicker:
        add_text(slide, Inches(3.2), H - FOOTER_H + Inches(0.08), Inches(7.4), Inches(0.30), kicker, 11, False, LINE)
    add_text(slide, Inches(11.2), H - FOOTER_H + Inches(0.08), Inches(1.9), Inches(0.30), f"{page}  /  {total}", 11, True, WHITE, PP_ALIGN.RIGHT)


def blank(prs, fill=WHITE):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, W, H, fill)
    return slide


def heading(slide, title, y=Inches(0.26), size=28):
    add_text(slide, Inches(0.5), y, Inches(12.3), Inches(0.55), title, size, True, TEAL_DARK)


def chip(slide, l, t, w, h, text, fill=TEAL, color=WHITE):
    add_round(slide, l, t, w, h, fill)
    add_text(slide, l, t + Inches(0.03), w, h - Inches(0.03), text, 11, True, color, PP_ALIGN.CENTER)


def card(slide, l, t, w, h, fill=CARD, accent=None):
    add_round(slide, l, t, w, h, fill)
    if accent:
        add_rect(slide, l, t, Inches(0.09), h, accent)


def title_slide(prs, title, subtitle, kicker, visual=None):
    s = blank(prs, CREAM)
    add_rect(s, 0, 0, Inches(0.16), H, TEAL)
    add_text(s, Inches(0.65), Inches(1.45), Inches(8.5), Inches(0.38), kicker, 15, True, TEAL_MID)
    add_text(s, Inches(0.65), Inches(1.95), Inches(8.7), Inches(2.1), title, 36, True, TEAL_DARK)
    add_text(s, Inches(0.65), Inches(4.25), Inches(8.5), Inches(1.3), subtitle, 17, False, MUTED)
    add_text(s, Inches(0.65), Inches(6.15), Inches(8.5), Inches(0.32), "PowerPoint teaching resource", 13, True, TEAL)
    if visual and Path(visual).exists():
        s.shapes.add_picture(str(visual), Inches(9.4), Inches(1.45), Inches(3.45), Inches(4.55))
    return s


def objectives_slide(prs, items):
    s = blank(prs)
    heading(s, "Learning objectives")
    add_text(s, Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.32), "By the end of this lesson you should be able to:", 15, False, MUTED)
    colours = [TEAL, TEAL_MID, NAVY, PURPLE, AMBER]
    n = len(items)
    h = Inches(0.92) if n <= 5 else Inches(0.78)
    for i, item in enumerate(items):
        y = Inches(1.32) + h * i + Inches(i * 0.08)
        add_round(s, Inches(0.5), y, Inches(12.3), h, CREAM)
        add_round(s, Inches(0.68), y + Inches(0.22), Inches(0.46), Inches(0.46), colours[i % len(colours)])
        add_text(s, Inches(0.68), y + Inches(0.26), Inches(0.46), Inches(0.38), str(i + 1), 15, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, Inches(1.35), y + Inches(0.22), Inches(11.2), h - Inches(0.3), item, 17, False, INK)
    return s


def section_slide(prs, title, subtitle="", points=None):
    s = blank(prs, TEAL_DARK)
    add_rect(s, 0, 0, Inches(0.16), H, TEAL)
    add_text(s, Inches(0.7), Inches(0.42), Inches(12.0), Inches(0.95), title, 30, True, WHITE)
    add_text(
        s,
        Inches(0.7),
        Inches(1.38),
        Inches(12.0),
        Inches(0.55),
        subtitle or "Move from a GCSE label to a Level 3 explanation.",
        17,
        False,
        LINE,
    )
    cards = points or [
        ("Look for", "The labelled diagram, symbol or equation that carries this idea."),
        ("Say in full", "A sentence that names the particles, structure or quantity — not ‘it’."),
        ("Level 3 move", subtitle or "Explain using attraction, structure or a calculation, not a one-word label."),
    ]
    accents = [TEAL, RGBColor(0x99, 0xF6, 0xE4), RGBColor(0x5E, 0xEA, 0xD4)]
    for i, (head, body) in enumerate(cards[:3]):
        x = Inches(0.7) + Inches(4.05) * i
        add_round(s, x, Inches(2.15), Inches(3.85), Inches(4.35), RGBColor(0x00, 0x3D, 0x33))
        add_rect(s, x, Inches(2.15), Inches(3.85), Inches(0.12), accents[i])
        add_text(s, x + Inches(0.22), Inches(2.42), Inches(3.4), Inches(0.5), f"{i+1}   {head}", 16, True, LINE)
        add_text(s, x + Inches(0.22), Inches(3.05), Inches(3.4), Inches(3.15), body, 16, False, WHITE)
    return s


def fact_cards(prs, title, cards):
    """2–4 equally sized cards that fill the teaching area."""
    s = blank(prs)
    heading(s, title)
    n = len(cards)
    gap = Inches(0.22)
    top = Inches(1.05)
    height = Inches(5.55)
    if n <= 3:
        width = (Inches(12.3) - gap * (n - 1)) / n
        colours = [TEAL, NAVY, PURPLE, AMBER][:n]
        for i, (head, body) in enumerate(cards):
            x = Inches(0.5) + (width + gap) * i
            add_round(s, x, top, width, height, CREAM)
            add_rect(s, x, top, width, Inches(0.12), colours[i])
            add_text(s, x + Inches(0.22), top + Inches(0.32), width - Inches(0.4), Inches(1.1), head, 18, True, TEAL_DARK)
            add_text(s, x + Inches(0.22), top + Inches(1.5), width - Inches(0.4), height - Inches(1.8), body, 15, False, INK)
    else:
        cols = 2
        rows = (n + 1) // 2
        width = (Inches(12.3) - gap) / 2
        ch = (height - gap * (rows - 1)) / rows
        colours = [TEAL, NAVY, PURPLE, AMBER, ROSE, SKY]
        for i, (head, body) in enumerate(cards):
            r, c = divmod(i, cols) if rows > 1 else (0, i)
            if n == 4:
                r, c = divmod(i, 2)
            x = Inches(0.5) + (width + gap) * c
            y = top + (ch + gap) * r
            add_round(s, x, y, width, ch, CREAM)
            add_rect(s, x, y, Inches(0.1), ch, colours[i % len(colours)])
            add_text(s, x + Inches(0.28), y + Inches(0.16), width - Inches(0.45), Inches(0.42), head, 16, True, TEAL_DARK)
            add_text(s, x + Inches(0.28), y + Inches(0.58), width - Inches(0.45), ch - Inches(0.74), body, 14, False, INK)
    return s


def process_steps(prs, title, steps, caption=""):
    s = blank(prs)
    heading(s, title)
    n = len(steps)
    gap = Inches(0.18)
    width = (Inches(12.3) - gap * (n - 1)) / n
    colours = [TEAL, TEAL_MID, NAVY, PURPLE, AMBER, ROSE]
    for i, (head, body) in enumerate(steps):
        x = Inches(0.5) + (width + gap) * i
        add_round(s, x, Inches(1.15), width, Inches(5.15), CREAM)
        add_round(s, x + width / 2 - Inches(0.28), Inches(1.35), Inches(0.56), Inches(0.56), colours[i % len(colours)])
        add_text(s, x + width / 2 - Inches(0.28), Inches(1.42), Inches(0.56), Inches(0.42), str(i + 1), 16, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.14), Inches(2.1), width - Inches(0.28), Inches(0.7), head, 15, True, TEAL_DARK, PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.16), Inches(2.85), width - Inches(0.32), Inches(3.2), body, 13, False, INK)
        if i < n - 1:
            add_text(s, x + width - Inches(0.08), Inches(3.4), Inches(0.28), Inches(0.4), "→", 20, True, TEAL)
    if caption:
        add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.32), caption, 13, False, MUTED)
    return s


def diagram_explain(prs, title, image_path, points, caption=""):
    """Labelled diagram on the left, short teaching cards on the right."""
    s = blank(prs)
    heading(s, title)
    if image_path and Path(image_path).exists():
        s.shapes.add_picture(str(image_path), Inches(0.4), Inches(1.05), Inches(6.7), Inches(5.15))
    colours = [TEAL, NAVY, PURPLE, AMBER, ROSE]
    n = max(len(points), 1)
    gap = Inches(0.1)
    usable = Inches(5.15)
    ch = (usable - gap * (n - 1)) / n
    if ch > Inches(1.22):
        ch = Inches(1.18)
    for i, point in enumerate(points):
        y = Inches(1.05) + (ch + gap) * i
        add_round(s, Inches(7.3), y, Inches(5.5), ch, CREAM)
        add_rect(s, Inches(7.3), y, Inches(0.1), ch, colours[i % len(colours)])
        add_text(s, Inches(7.58), y + Inches(0.1), Inches(5.05), ch - Inches(0.16), point, 14, False, INK)
    leftover_top = Inches(1.05) + (ch + gap) * n
    if leftover_top < Inches(6.05):
        add_round(s, Inches(7.3), leftover_top, Inches(5.5), Inches(6.35) - leftover_top, NAVY_BG)
        add_text(
            s,
            Inches(7.5),
            leftover_top + Inches(0.1),
            Inches(5.15),
            Inches(6.25) - leftover_top,
            caption or "Say the labelled feature and its function in one sentence.",
            13,
            False,
            NAVY,
        )
    elif caption:
        add_text(s, Inches(0.4), Inches(6.28), Inches(6.7), Inches(0.32), caption, 12, False, MUTED)
    return s


def worked_example(prs, title, question, equation, substitution, working, answer, unit, note=""):
    s = blank(prs, SKY_BG)
    chip(s, Inches(0.5), Inches(0.22), Inches(2.6), Inches(0.34), "WORKED EXAMPLE", SKY)
    heading(s, title, Inches(0.62), 26)
    boxes = [
        ("Question", question, NAVY, NAVY_BG, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.95)),
        ("Equation", equation, TEAL, CREAM, Inches(0.5), Inches(2.38), Inches(6.0), Inches(1.15)),
        ("Substitution", substitution, PURPLE, PURPLE_BG, Inches(6.7), Inches(2.38), Inches(6.1), Inches(1.15)),
        ("Working", working, AMBER, AMBER_BG, Inches(0.5), Inches(3.66), Inches(6.0), Inches(1.45)),
        ("Answer + unit", f"{answer}     {unit}", GREEN, GREEN_BG, Inches(6.7), Inches(3.66), Inches(6.1), Inches(1.45)),
    ]
    for label, body, accent, fill, x, y, w, h in boxes:
        add_round(s, x, y, w, h, fill)
        add_rect(s, x, y, Inches(0.1), h, accent)
        add_text(s, x + Inches(0.24), y + Inches(0.08), w - Inches(0.4), Inches(0.28), label, 12, True, accent)
        add_text(s, x + Inches(0.24), y + Inches(0.38), w - Inches(0.4), h - Inches(0.48), body, 16, False, INK)
    if note:
        add_text(s, Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.35), note, 15, False, MUTED)
    else:
        add_round(s, Inches(0.5), Inches(5.25), Inches(12.3), Inches(1.35), WHITE)
        add_text(s, Inches(0.74), Inches(5.4), Inches(11.8), Inches(1.05), "Teacher note: write the equation first, substitute with units, then calculate. Sense-check the magnitude.", 15, False, MUTED)
    return s


def question_cards(prs, title, items, kind="Check your understanding"):
    """items: list of (text, marks) or strings."""
    s = blank(prs, ROSE_BG)
    chip(s, Inches(0.5), Inches(0.22), Inches(3.6), Inches(0.34), kind.upper(), ROSE)
    heading(s, title, Inches(0.64), 26)
    n = len(items)
    gap = Inches(0.16)
    height = Inches(4.85)
    width = (Inches(12.3) - gap * (n - 1)) / n if n else Inches(12.3)
    if n <= 3:
        for i, item in enumerate(items):
            text, marks = item if isinstance(item, tuple) else (item, "")
            x = Inches(0.5) + (width + gap) * i
            add_round(s, x, Inches(1.3), width, height, WHITE)
            add_round(s, x + Inches(0.2), Inches(1.5), Inches(0.5), Inches(0.5), ROSE)
            add_text(s, x + Inches(0.2), Inches(1.56), Inches(0.5), Inches(0.4), str(i + 1), 16, True, WHITE, PP_ALIGN.CENTER)
            if marks:
                add_text(s, x + Inches(0.8), Inches(1.56), width - Inches(1.1), Inches(0.4), f"{marks} marks", 13, True, ROSE)
            add_text(s, x + Inches(0.22), Inches(2.2), width - Inches(0.44), Inches(1.55), text, 16, False, INK)
            add_text(s, x + Inches(0.22), Inches(3.85), width - Inches(0.44), Inches(0.28), "Write a full-sentence answer", 11, True, MUTED)
            for k in range(4):
                add_rect(s, x + Inches(0.22), Inches(4.22) + Inches(0.32) * k, width - Inches(0.44), Pt(1.25), RGBColor(0xE2, 0xE8, 0xF0))
    else:
        rows = 2
        cols = (n + 1) // 2
        cw = (Inches(12.3) - gap * (cols - 1)) / cols
        ch = (height - gap) / rows
        for i, item in enumerate(items):
            text, marks = item if isinstance(item, tuple) else (item, "")
            r, c = divmod(i, cols)
            x = Inches(0.5) + (cw + gap) * c
            y = Inches(1.3) + (ch + gap) * r
            add_round(s, x, y, cw, ch, WHITE)
            add_text(s, x + Inches(0.18), y + Inches(0.12), Inches(0.4), Inches(0.32), f"{i+1}.", 14, True, ROSE)
            add_text(s, x + Inches(0.5), y + Inches(0.12), cw - Inches(0.7), ch - Inches(0.28), text + (f"  [{marks}]" if marks else ""), 14, False, INK)
    add_text(s, Inches(0.5), Inches(6.28), Inches(12.3), Inches(0.3), "Think first — answers are on the next slide.", 13, True, ROSE)
    return s


def answer_cards(prs, title, items):
    s = blank(prs, GREEN_BG)
    chip(s, Inches(0.5), Inches(0.22), Inches(2.2), Inches(0.34), "ANSWERS", GREEN)
    heading(s, title, Inches(0.64), 26)
    n = len(items)
    gap = Inches(0.16)
    height = Inches(5.15)
    width = (Inches(12.3) - gap * (n - 1)) / n if n <= 3 else Inches(12.3)
    if n <= 3:
        for i, item in enumerate(items):
            x = Inches(0.5) + (width + gap) * i
            add_round(s, x, Inches(1.3), width, height, WHITE)
            add_round(s, x + Inches(0.2), Inches(1.5), Inches(0.5), Inches(0.5), GREEN)
            add_text(s, x + Inches(0.2), Inches(1.56), Inches(0.5), Inches(0.4), str(i + 1), 16, True, WHITE, PP_ALIGN.CENTER)
            add_text(s, x + Inches(0.22), Inches(2.2), width - Inches(0.44), height - Inches(1.15), item, 15, False, INK)
    else:
        ch = height / n - Inches(0.06)
        for i, item in enumerate(items):
            y = Inches(1.25) + (ch + Inches(0.08)) * i
            add_round(s, Inches(0.5), y, Inches(12.3), ch, WHITE)
            add_round(s, Inches(0.68), y + Inches(0.12), Inches(0.4), Inches(0.4), GREEN)
            add_text(s, Inches(0.68), y + Inches(0.16), Inches(0.4), Inches(0.32), str(i + 1), 13, True, WHITE, PP_ALIGN.CENTER)
            add_text(s, Inches(1.25), y + Inches(0.12), Inches(11.3), ch - Inches(0.18), item, 15, False, INK)
    return s


def two_col(prs, title, left_title, left, right_title, right, left_fill=CREAM, right_fill=NAVY_BG, left_accent=TEAL, right_accent=NAVY):
    s = blank(prs)
    heading(s, title)
    card(s, Inches(0.5), Inches(1.1), Inches(6.05), Inches(5.5), left_fill, left_accent)
    add_text(s, Inches(0.8), Inches(1.28), Inches(5.5), Inches(0.4), left_title, 18, True, left_accent)
    add_bullets(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.55), left, 15)
    card(s, Inches(6.75), Inches(1.1), Inches(6.05), Inches(5.5), right_fill, right_accent)
    add_text(s, Inches(7.05), Inches(1.28), Inches(5.5), Inches(0.4), right_title, 18, True, right_accent)
    add_bullets(s, Inches(7.05), Inches(1.8), Inches(5.5), Inches(4.55), right, 15)
    return s


def table_slide(prs, title, headers, rows, note=""):
    s = blank(prs)
    heading(s, title)
    cols = len(headers)
    table_w = Inches(12.3)
    left = Inches(0.5)
    top = Inches(1.1)
    usable = Inches(5.4) if note else Inches(5.7)
    row_h = min(Inches(0.72), usable / (len(rows) + 1))
    add_rect(s, left, top, table_w, row_h, TEAL_DARK)
    col_w = table_w // cols
    for i, h in enumerate(headers):
        add_text(s, left + Emu(i * col_w), top + Inches(0.12), Emu(col_w), row_h - Inches(0.16), h, 13, True, WHITE, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = top + row_h * (r + 1)
        add_rect(s, left, y, table_w, row_h, CREAM if r % 2 == 0 else CARD)
        for i, cell in enumerate(row):
            add_text(s, left + Emu(i * col_w) + Inches(0.08), y + Inches(0.1), Emu(col_w) - Inches(0.12), row_h - Inches(0.14), cell, 13, False, INK, PP_ALIGN.CENTER)
    if note:
        add_text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.3), note, 13, False, MUTED)
    return s


def misconception_slide(prs, pairs):
    s = blank(prs)
    heading(s, "Common misconceptions")
    n = len(pairs)
    block = Inches(5.5) / n - Inches(0.08)
    for i, (wrong, right) in enumerate(pairs):
        y = Inches(1.1) + (block + Inches(0.12)) * i
        card(s, Inches(0.5), y, Inches(6.05), block, ROSE_BG, ROSE)
        add_text(s, Inches(0.75), y + Inches(0.1), Inches(5.55), Inches(0.26), "Not quite", 12, True, ROSE)
        add_text(s, Inches(0.75), y + Inches(0.38), Inches(5.55), block - Inches(0.5), wrong, 14, False, INK)
        card(s, Inches(6.75), y, Inches(6.05), block, GREEN_BG, GREEN)
        add_text(s, Inches(7.0), y + Inches(0.1), Inches(5.55), Inches(0.26), "Instead", 12, True, GREEN)
        add_text(s, Inches(7.0), y + Inches(0.38), Inches(5.55), block - Inches(0.5), right, 14, False, INK)
    return s


def activity_slide(prs, title, steps, minutes="8 minutes"):
    s = blank(prs, AMBER_BG)
    chip(s, Inches(0.5), Inches(0.22), Inches(3.0), Inches(0.34), "STUDENT ACTIVITY", AMBER)
    heading(s, title, Inches(0.64), 26)
    add_text(s, Inches(10.0), Inches(0.68), Inches(2.8), Inches(0.32), minutes, 14, True, AMBER, PP_ALIGN.RIGHT)
    n = len(steps)
    ch = Inches(5.05) / n - Inches(0.08)
    for i, step in enumerate(steps):
        y = Inches(1.25) + (ch + Inches(0.1)) * i
        add_round(s, Inches(0.5), y, Inches(12.3), ch, WHITE)
        add_round(s, Inches(0.7), y + Inches(0.16), Inches(0.46), Inches(0.46), AMBER)
        add_text(s, Inches(0.7), y + Inches(0.2), Inches(0.46), Inches(0.38), str(i + 1), 14, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, Inches(1.4), y + Inches(0.16), Inches(11.15), ch - Inches(0.24), step, 16, False, INK)
    return s


def whiteboard_slide(prs, title, prompt, boxes):
    s = blank(prs)
    chip(s, Inches(0.5), Inches(0.22), Inches(3.3), Inches(0.34), "MINI-WHITEBOARD", NAVY)
    heading(s, title, Inches(0.64), 26)
    add_text(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.55), prompt, 16, False, MUTED)
    n = len(boxes)
    width = (Inches(12.3) - Inches(0.2) * (n - 1)) / n
    for i, box in enumerate(boxes):
        x = Inches(0.5) + (width + Inches(0.2)) * i
        add_round(s, x, Inches(1.95), width, Inches(4.5), RGBColor(0xFF, 0xFF, 0xF4), RGBColor(0xCB, 0xD5, 0xE1))
        add_text(s, x + Inches(0.18), Inches(2.1), width - Inches(0.36), Inches(0.4), box, 15, True, NAVY)
        add_text(s, x + Inches(0.18), Inches(5.9), width - Inches(0.36), Inches(0.35), "Write, then hold up.", 12, False, MUTED)
    return s


def match_slide(prs, title, left, right, instruction="Draw lines to match. Then justify one pair."):
    s = blank(prs, PURPLE_BG)
    chip(s, Inches(0.5), Inches(0.22), Inches(3.4), Inches(0.34), "MATCHING TASK", PURPLE)
    heading(s, title, Inches(0.64), 26)
    add_text(s, Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.35), instruction, 14, False, MUTED)
    n = len(left)
    ch = Inches(4.85) / n - Inches(0.08)
    for i, (a, b) in enumerate(zip(left, right)):
        y = Inches(1.7) + (ch + Inches(0.1)) * i
        add_round(s, Inches(0.5), y, Inches(5.3), ch, WHITE)
        add_text(s, Inches(0.7), y + Inches(0.12), Inches(4.95), ch - Inches(0.2), a, 15, True, TEAL_DARK)
        add_round(s, Inches(7.5), y, Inches(5.3), ch, WHITE)
        add_text(s, Inches(7.7), y + Inches(0.12), Inches(4.95), ch - Inches(0.2), b, 15, False, INK)
        add_text(s, Inches(6.0), y + Inches(0.1), Inches(1.3), ch - Inches(0.15), "— ? —", 14, True, PURPLE, PP_ALIGN.CENTER)
    return s


def calc_scaffold(prs, title, given, find, hint=""):
    s = blank(prs)
    heading(s, title)
    card(s, Inches(0.5), Inches(1.1), Inches(6.05), Inches(2.3), CREAM, TEAL)
    add_text(s, Inches(0.8), Inches(1.25), Inches(5.5), Inches(0.35), "Given", 14, True, TEAL)
    add_bullets(s, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.5), given, 15)
    card(s, Inches(6.75), Inches(1.1), Inches(6.05), Inches(2.3), NAVY_BG, NAVY)
    add_text(s, Inches(7.05), Inches(1.25), Inches(5.5), Inches(0.35), "Find", 14, True, NAVY)
    add_bullets(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(1.5), find, 15)
    labels = ["Equation", "Substitution", "Working", "Answer + unit"]
    for i, lab in enumerate(labels):
        x = Inches(0.5) + Inches(3.15) * i
        add_round(s, x, Inches(3.6), Inches(3.0), Inches(2.95), CARD)
        add_text(s, x + Inches(0.15), Inches(3.75), Inches(2.7), Inches(0.35), lab, 14, True, TEAL_DARK)
        add_text(s, x + Inches(0.15), Inches(4.2), Inches(2.7), Inches(2.15), " ", 14)
    if hint:
        add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.22), hint, 12, False, MUTED)
    return s


def plenary_slide(prs, items):
    s = blank(prs)
    heading(s, "Plenary")
    add_text(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.35), "30–60 seconds each. No notes if you can avoid them.", 14, False, MUTED)
    n = len(items)
    width = (Inches(12.3) - Inches(0.2) * (n - 1)) / n
    colours = [TEAL, NAVY, PURPLE, AMBER]
    for i, item in enumerate(items):
        x = Inches(0.5) + (width + Inches(0.2)) * i
        add_round(s, x, Inches(1.4), width, Inches(5.15), CREAM)
        add_round(s, x + Inches(0.2), Inches(1.6), Inches(0.5), Inches(0.5), colours[i % 4])
        add_text(s, x + Inches(0.2), Inches(1.68), Inches(0.5), Inches(0.38), str(i + 1), 16, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), Inches(2.3), width - Inches(0.4), Inches(4.0), item, 16, False, INK)
    return s


def homework_slide(prs, items, next_lesson=""):
    s = blank(prs)
    heading(s, "Independent practice")
    for i, item in enumerate(items):
        y = Inches(1.15) + Inches(1.15) * i
        add_round(s, Inches(0.5), y, Inches(12.3), Inches(1.02), CREAM)
        add_round(s, Inches(0.7), y + Inches(0.26), Inches(0.5), Inches(0.5), TEAL)
        add_text(s, Inches(0.7), y + Inches(0.32), Inches(0.5), Inches(0.4), str(i + 1), 15, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, Inches(1.45), y + Inches(0.28), Inches(11.1), Inches(0.55), item, 16, False, INK)
    if next_lesson:
        add_text(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.3), next_lesson, 14, True, TEAL)
    return s


def depth_check(prs, rows):
    """Internal Level 3 content-depth reminder for teachers."""
    s = blank(prs, NAVY_BG)
    chip(s, Inches(0.5), Inches(0.22), Inches(3.6), Inches(0.34), "LEVEL 3 DEPTH CHECK", NAVY)
    heading(s, "Beyond GCSE recall", Inches(0.64), 26)
    add_text(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4), "This lesson must progress past simple recall. Use these Level 3 moves in explanations and answers.", 14, False, MUTED)
    table_slide_on(s, ["GCSE starting point", "Unit 1 progression"], rows, Inches(1.75))
    return s


def table_slide_on(s, headers, rows, top):
    cols = len(headers)
    table_w = Inches(12.3)
    left = Inches(0.5)
    usable = Inches(5.0)
    row_h = min(Inches(0.78), usable / (len(rows) + 1))
    add_rect(s, left, top, table_w, row_h, TEAL_DARK)
    col_w = table_w // cols
    for i, h in enumerate(headers):
        add_text(s, left + Emu(i * col_w), top + Inches(0.14), Emu(col_w), row_h - Inches(0.18), h, 13, True, WHITE, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = top + row_h * (r + 1)
        add_rect(s, left, y, table_w, row_h, WHITE if r % 2 else CREAM)
        for i, cell in enumerate(row):
            add_text(s, left + Emu(i * col_w) + Inches(0.12), y + Inches(0.1), Emu(col_w) - Inches(0.2), row_h - Inches(0.16), cell, 13, False, INK)


def apply_footers(prs, kicker: str):
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        footer(slide, i, total, kicker)


def set_core(prs, title: str, subject: str):
    prs.core_properties.title = title
    prs.core_properties.author = "JDScience"
    prs.core_properties.subject = subject
    prs.core_properties.keywords = subject
    prs.core_properties.category = "Teaching resource"


def save_prs(prs, path: Path, title: str, subject: str, kicker: str):
    set_core(prs, title, subject)
    apply_footers(prs, kicker)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


# Compatibility aliases used by older lesson modules during the rebuild.
def content_slide(prs, title, bullets, aside=None, aside_title="Remember"):
    cards = []
    for i, b in enumerate(bullets[:4]):
        cards.append((f"Point {i+1}" if len(b) > 80 else b.split(".")[0][:40], b))
    if aside:
        cards = [(bullets[0][:40], bullets[0])]
        return two_col(prs, title, title, bullets, aside_title, aside)
    return fact_cards(prs, title, [(f"{i+1}", b) for i, b in enumerate(bullets[:4])])


def question_slide(prs, title, questions, kind="Check your understanding"):
    return question_cards(prs, title, questions, kind)


def answer_slide(prs, title, answers):
    return answer_cards(prs, title, answers)


def diagram_slide(prs, title, image_path, caption="", bullets=None):
    return diagram_explain(prs, title, image_path, bullets or [], caption)
