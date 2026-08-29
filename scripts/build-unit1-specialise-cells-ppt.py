#!/usr/bin/env python3
"""Build the JDScience-branded Unit 1 specialise cells PowerPoint."""

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
MEDIA = ROOT / "content/shop/unit-1-specialise-cells/media"
OUT_PPT = ROOT / "content/shop/Unit 1 specialise cells.pptx"
OUT_COVER = ROOT / "content/shop/unit-1-specialise-cells/cover.png"

W, H = Inches(13.333), Inches(7.5)
TEAL = RGBColor(0x00, 0x96, 0x88)
TEAL_DARK = RGBColor(0x00, 0x4D, 0x40)
TEAL_MID = RGBColor(0x00, 0x79, 0x6B)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF0, 0xFD, 0xFA)
CARD = RGBColor(0xF8, 0xFA, 0xFC)
LINE = RGBColor(0xCC, 0xFB, 0xF1)

TITLE = "Unit 1 specialise cells"


def set_run(run, size=18, bold=False, color=INK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_textbox(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_round(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def brand_chrome(slide, n, total, section=""):
    add_rect(slide, 0, 0, W, Inches(0.72), TEAL_DARK)
    add_rect(slide, 0, Inches(0.72), W, Inches(0.08), TEAL)
    add_textbox(slide, Inches(0.35), Inches(0.12), Inches(3.2), Inches(0.5), "JD SCIENCE", 20, True, WHITE)
    add_textbox(
        slide,
        Inches(4.0),
        Inches(0.16),
        Inches(8.9),
        Inches(0.42),
        section or TITLE,
        14,
        False,
        RGBColor(0xCC, 0xFB, 0xF1),
        PP_ALIGN.RIGHT,
    )
    add_rect(slide, 0, Inches(7.18), W, Inches(0.32), TEAL_DARK)
    add_textbox(slide, Inches(0.35), Inches(7.18), Inches(8.5), Inches(0.3), "jdscience.co.uk  ·  Specialised Cells", 11, False, WHITE)
    add_textbox(slide, Inches(10.6), Inches(7.18), Inches(2.3), Inches(0.3), f"{n} / {total}", 11, False, WHITE, PP_ALIGN.RIGHT)


def bullets(slide, items, l=Inches(0.45), t=Inches(1.15), w=Inches(12.4), h=Inches(5.8), size=22):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=INK)
    return box


def picture(slide, name, l, t, w, h):
    path = MEDIA / name
    return slide.shapes.add_picture(str(path), l, t, w, h)


def heading(slide, text, t=Inches(1.0), size=30):
    add_textbox(slide, Inches(0.45), t, Inches(12.4), Inches(0.6), text, size, True, TEAL_DARK)


def card(slide, l, t, w, h):
    shape = add_round(slide, l, t, w, h, CARD)
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    return shape


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    slides = []

    def new(section=""):
        s = prs.slides.add_slide(blank)
        slides.append(s)
        return s

    # 1 Title
    s = new()
    add_rect(s, 0, 0, W, H, TEAL_DARK)
    add_rect(s, 0, 0, Inches(0.22), H, TEAL)
    add_textbox(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4), "JD SCIENCE  ·  BIOLOGY", 16, True, TEAL)
    add_textbox(s, Inches(0.7), Inches(2.15), Inches(12), Inches(1.3), TITLE, 48, True, WHITE)
    add_textbox(s, Inches(0.7), Inches(3.5), Inches(12), Inches(0.6), "Specialised Cells", 28, False, RGBColor(0xCC, 0xFB, 0xF1))
    add_textbox(
        s,
        Inches(0.7),
        Inches(4.4),
        Inches(11),
        Inches(1.6),
        "How cell structure is adapted to function — palisade mesophyll cells, root hair cells, sperm cells, egg cells, red blood cells and neutrophils. Full diagrams included.",
        18,
        False,
        WHITE,
    )
    add_textbox(s, Inches(0.7), Inches(6.6), Inches(11), Inches(0.4), "jdscience.co.uk", 16, True, TEAL)

    # 2 TEM eukaryotic cell
    s = new("Recap · eukaryotic ultrastructure")
    brand_chrome(s, 2, 29, "Recap")
    heading(s, "Eukaryotic cell ultrastructure")
    picture(s, "image3.png", Inches(1.6), Inches(1.65), Inches(10.1), Inches(4.85))
    add_textbox(
        s,
        Inches(0.45),
        Inches(6.55),
        Inches(12.4),
        Inches(0.5),
        "Transmission electron micrograph of a eukaryotic cell. Scale bar: 2 μm. Note the nucleus, mitochondria with cristae, and endomembrane system.",
        13,
        False,
        MUTED,
    )

    # 3 Recap plant cell TEM
    s = new("Recap")
    brand_chrome(s, 3, 29, "Recap")
    heading(s, "Recap")
    picture(s, "image4.png", Inches(2.4), Inches(1.6), Inches(8.5), Inches(5.4))
    add_textbox(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.5), "Labels on the micrograph: N = nucleus  ·  M = mitochondria", 13, False, MUTED)

    # 4 Recap bacteria TEM
    s = new("Recap")
    brand_chrome(s, 4, 29, "Recap")
    heading(s, "Recap")
    picture(s, "image5.png", Inches(2.55), Inches(1.6), Inches(8.2), Inches(5.4))
    add_textbox(
        s,
        Inches(0.45),
        Inches(6.55),
        Inches(12.4),
        Inches(0.5),
        "Rod-shaped bacteria (prokaryotes). The central cell is undergoing binary fission. No membrane-bound organelles are present.",
        13,
        False,
        MUTED,
    )

    # 5 Starter
    s = new("Starter")
    brand_chrome(s, 5, 29, "Starter")
    heading(s, "What is a specialised cell?", size=34)
    card(s, Inches(0.45), Inches(1.9), Inches(12.4), Inches(2.0))
    add_textbox(s, Inches(0.75), Inches(2.15), Inches(11.8), Inches(1.5), "What is a specialised cell?", 32, True, TEAL_DARK)
    card(s, Inches(0.45), Inches(4.2), Inches(12.4), Inches(2.0))
    add_textbox(s, Inches(0.75), Inches(4.45), Inches(11.8), Inches(1.5), "Do you have any examples?", 32, True, TEAL_DARK)

    # 6 Plant cell diagram
    s = new("Plant cell")
    brand_chrome(s, 6, 29, "Plant cell")
    heading(s, "Plant cell — labelled diagram")
    picture(s, "image6.png", Inches(2.4), Inches(1.55), Inches(8.5), Inches(5.4))

    # 7 Palisade 1
    s = new("Palisade mesophyll cell")
    brand_chrome(s, 7, 29, "Palisade mesophyll cell")
    heading(s, "Palisade Mesophyll Cell")
    bullets(
        s,
        [
            "Palisade mesophyll cells found in leaves",
            "They are rectangular box-shaped cells that contain chloroplasts",
            "The chloroplasts can absorb a large amount of light for photosynthesis",
        ],
    )

    # 8 Palisade 2
    s = new("Palisade mesophyll cell")
    brand_chrome(s, 8, 29, "Palisade mesophyll cell")
    heading(s, "Palisade Mesophyll Cell")
    bullets(
        s,
        [
            "They also move around in the cytoplasm to maximise amount of light absorbed",
            "The cells are closely packed to form a continuous layer in the leaf",
            "Palisade cells are surrounded by a plasma membrane and a cell wall made of cellulose",
            "This protects the cell and keeps it rigid",
        ],
    )

    # 9 Palisade 3
    s = new("Palisade mesophyll cell")
    brand_chrome(s, 9, 29, "Palisade mesophyll cell")
    heading(s, "Palisade Mesophyll Cell")
    bullets(
        s,
        [
            "They also have a large vacuole to maintain turgor pressure",
            "This means the plasma membrane pushes against the cell wall of the plant to maintain its rigid structure",
        ],
    )

    # 10 Root hair diagram
    s = new("Root hair cell")
    brand_chrome(s, 10, 29, "Root hair cell")
    heading(s, "Root hair cell — labelled diagram")
    picture(s, "image7.png", Inches(1.3), Inches(1.55), Inches(10.7), Inches(5.4))

    # 11 Root hair 1
    s = new("Root hair cell")
    brand_chrome(s, 11, 29, "Root hair cell")
    heading(s, "Root hair cell")
    bullets(
        s,
        [
            "These are found at the plants roots, near the growing tip",
            "They have long hair-like extensions called root hairs",
            "The root hairs increase the surface area of the cell to maximise the movement of the water and minerals from the soil into the plant root",
        ],
    )

    # 12 Root hair 2
    s = new("Root hair cell")
    brand_chrome(s, 12, 29, "Root hair cell")
    heading(s, "Root hair cell")
    bullets(
        s,
        [
            "The cells have thin cellulose walls and a vacuole containing cell sap with a low water potential",
            "This encourages movement of water into the cell",
        ],
    )

    # 13 Sperm diagram
    s = new("Sperm cell")
    brand_chrome(s, 13, 29, "Sperm cell")
    heading(s, "Sperm cell — labelled diagram")
    picture(s, "image8.png", Inches(1.15), Inches(1.55), Inches(11.0), Inches(5.4))

    # 14 Sperm labels
    s = new("Sperm cell")
    brand_chrome(s, 14, 29, "Sperm cell")
    heading(s, "Sperm cell — key structures")
    labels = [
        ("Undulipodium", "Tail-like structure used for movement"),
        ("Plasma membrane", "Outer boundary of the cell"),
        ("Acrosome", "Cap containing digestive enzymes"),
        ("Nucleus", "Contains the haploid genetic information"),
        ("Mitochondria", "Supply the energy needed for movement"),
    ]
    for i, (title, detail) in enumerate(labels):
        col = i % 3
        row = i // 3
        l = Inches(0.45 + col * 4.2)
        t = Inches(1.85 + row * 2.35)
        card(s, l, t, Inches(4.0), Inches(2.15))
        add_textbox(s, l + Inches(0.2), t + Inches(0.25), Inches(3.6), Inches(0.5), title, 20, True, TEAL_DARK)
        add_textbox(s, l + Inches(0.2), t + Inches(0.85), Inches(3.6), Inches(1.05), detail, 16, False, MUTED)

    # 15 Sperm 1
    s = new("Sperm cell")
    brand_chrome(s, 15, 29, "Sperm cell")
    heading(s, "Sperm Cell")
    bullets(
        s,
        [
            "These are male gametes",
            "They have a tail like structure called a undulipodium so they can move",
            "They also contain many mitochondria to supply the energy needed for this movement",
        ],
    )

    # 16 Sperm 2
    s = new("Sperm cell")
    brand_chrome(s, 16, 29, "Sperm cell")
    heading(s, "Sperm Cell")
    bullets(
        s,
        [
            "It is made up of an acrosome which contains digestive enzymes",
            "These enzymes are released when the sperm meets the egg",
            "This digests the protective layer and allows the sperm to penetrate",
        ],
    )

    # 17 Sperm 3
    s = new("Sperm cell")
    brand_chrome(s, 17, 29, "Sperm cell")
    heading(s, "Sperm Cell")
    card(s, Inches(0.45), Inches(2.0), Inches(12.4), Inches(3.4))
    add_textbox(
        s,
        Inches(0.8),
        Inches(2.4),
        Inches(11.7),
        Inches(2.6),
        "The sperms function is to deliver genetic information to the egg cell or ovum (male gamete).",
        28,
        True,
        TEAL_DARK,
    )

    # 18 Egg and sperm comparison
    s = new("Egg cell")
    brand_chrome(s, 18, 29, "Egg cell and sperm cell")
    heading(s, "Egg and sperm — comparison diagrams")
    picture(s, "image9.png", Inches(1.35), Inches(1.55), Inches(10.6), Inches(5.4))

    # 19 Egg labels
    s = new("Egg cell")
    brand_chrome(s, 19, 29, "Egg cell")
    heading(s, "Egg cell — key structures")
    egg_labels = [
        ("Corona radiata", "Two or three outer layers attached to the zona pellucida. Supplies proteins needed to the fertilised egg."),
        ("Zona pellucida", "The outer protective layer/membrane of the egg."),
        ("Nucleus", "Contains the haploid genetic information."),
        ("Cytoplasm", "Contains nutrients and organelles needed after fertilisation."),
        ("Cell membrane", "Plasma membrane immediately around the cytoplasm."),
    ]
    for i, (title, detail) in enumerate(egg_labels):
        col = i % 3
        row = i // 3
        l = Inches(0.45 + col * 4.2)
        t = Inches(1.85 + row * 2.35)
        card(s, l, t, Inches(4.0), Inches(2.15))
        add_textbox(s, l + Inches(0.2), t + Inches(0.2), Inches(3.6), Inches(0.45), title, 18, True, TEAL_DARK)
        add_textbox(s, l + Inches(0.2), t + Inches(0.7), Inches(3.6), Inches(1.25), detail, 14, False, MUTED)

    # 20 Egg intro
    s = new("Egg cell")
    brand_chrome(s, 20, 29, "Egg cell")
    heading(s, "Egg Cell")
    bullets(
        s,
        [
            "Egg cells or ova are the female gametes",
            "AN egg cell is one of the largest cells in the human body",
            "It is 0.12mm in diameter",
        ],
    )

    # 21 Egg detail
    s = new("Egg cell")
    brand_chrome(s, 21, 29, "Egg cell")
    heading(s, "Egg Cell")
    bullets(
        s,
        [
            "It has a nucleus",
            "A zona pellucida which is the outer protective layer/membrane of the egg",
            "Attached to this is the corona radiata which consists of two or three layers",
            "This functions to supply proteins needed to the fertilised egg.",
        ],
    )

    # 22 RBC image
    s = new("Red blood cell")
    brand_chrome(s, 22, 29, "Red blood cell")
    heading(s, "Red blood cell")
    picture(s, "image10.png", Inches(2.7), Inches(1.55), Inches(7.9), Inches(5.4))

    # 23 RBC 1
    s = new("Red blood cell")
    brand_chrome(s, 23, 29, "Red blood cell")
    heading(s, "Red blood cell")
    bullets(
        s,
        [
            "Red blood cells or erythrocytes are a biconcave disc",
            "This increases the surface area o volume ratio of an erythrocyte",
            "They are flexible so they can squeeze through narrow capillaries",
        ],
    )

    # 24 RBC 2
    s = new("Red blood cell")
    brand_chrome(s, 24, 29, "Red blood cell")
    heading(s, "Red blood cell")
    bullets(
        s,
        [
            "Their function is to transport oxygen around the body",
            "They do not have a nucleus",
            "This is to increase space for haemoglobin molecules to attach oxygen",
        ],
    )

    # 25 Neutrophil diagram
    s = new("White blood cell")
    brand_chrome(s, 25, 29, "White blood cell")
    heading(s, "Neutrophil — labelled diagram")
    picture(s, "image11.png", Inches(1.55), Inches(1.55), Inches(10.2), Inches(5.4))

    # 26 Neutrophil 1
    s = new("White blood cell")
    brand_chrome(s, 26, 29, "White blood cell")
    heading(s, "White blood cell")
    bullets(
        s,
        [
            "Neutrophils are types of white eblood cells that play an important part in the immune system",
            "They have multi lobed nuclei, enabling them to squeeze through small gaps when travelling to the site of an infection",
        ],
    )

    # 27 Neutrophil 2
    s = new("White blood cell")
    brand_chrome(s, 27, 29, "White blood cell")
    heading(s, "White blood cell")
    bullets(
        s,
        [
            "The cytoplasm holds lysosomes that contain enzymes that are used to digest pathogens that are ingested by the neutrophil",
        ],
    )

    # 28 Who am I
    s = new("Plenary")
    brand_chrome(s, 28, 29, "Plenary")
    heading(s, "Who am I?", size=36)
    card(s, Inches(0.45), Inches(1.95), Inches(12.4), Inches(4.4))
    add_textbox(
        s,
        Inches(0.85),
        Inches(2.3),
        Inches(11.6),
        Inches(3.6),
        "Describe the key features of different specialised cells – students guess the cell type",
        28,
        True,
        TEAL_DARK,
    )

    # 29 Homework
    s = new("Homework")
    brand_chrome(s, 29, 29, "Homework")
    heading(s, "Homework")
    card(s, Inches(0.45), Inches(1.95), Inches(12.4), Inches(4.4))
    add_textbox(s, Inches(0.85), Inches(2.4), Inches(11.6), Inches(1.0), "Revise for exam next week", 34, True, TEAL_DARK)
    add_textbox(s, Inches(0.85), Inches(3.6), Inches(11.6), Inches(2.0), "Review every specialised cell in this pack: palisade mesophyll, root hair, sperm, egg, red blood cell and neutrophil. Learn structure, adaptations and function.", 20, False, MUTED)

    core = prs.core_properties
    core.title = TITLE
    core.author = "JD Science"
    core.subject = "Specialised Cells"
    core.keywords = "specialised cells, biology, T Level, JDScience"

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"wrote {OUT_PPT} slides={len(prs.slides)}")

    # Shop cover
    cover = Image.new("RGB", (1600, 900), (0, 77, 64))
    draw = ImageDraw.Draw(cover)
    draw.rectangle((0, 0, 28, 900), fill=(0, 150, 136))
    try:
        font_lg = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
        font_md = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 36)
        font_sm = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
    except OSError:
        font_lg = font_md = font_sm = ImageFont.load_default()
    draw.text((80, 180), "JD SCIENCE", fill=(0, 150, 136), font=font_sm)
    draw.text((80, 250), TITLE, fill=(255, 255, 255), font=font_lg)
    draw.text((80, 360), "Specialised Cells  ·  PowerPoint", fill=(204, 251, 241), font=font_md)
    draw.text((80, 760), "jdscience.co.uk", fill=(0, 150, 136), font=font_sm)
    # small cell montage
    try:
        rbc = Image.open(MEDIA / "image10.png").convert("RGB")
        rbc.thumbnail((420, 420))
        cover.paste(rbc, (1100, 240))
    except Exception:
        pass
    cover.save(OUT_COVER, "PNG")
    print(f"wrote {OUT_COVER}")


if __name__ == "__main__":
    build()
